# Copyright (c) 2026 Beijing Volcano Engine Technology Co., Ltd.
# SPDX-License-Identifier: Apache-2.0
"""
PowerPoint (.pptx) parser for OpenViking.

Converts PowerPoint presentations to Markdown then parses using MarkdownParser.
Inspired by microsoft/markitdown approach.
"""

from pathlib import Path
from typing import List, Optional, Union

from openviking.parse.base import ParseResult
from openviking.parse.parsers.base_parser import BaseParser
from openviking_cli.utils.config.parser_config import ParserConfig
from openviking_cli.utils.logger import get_logger

logger = get_logger(__name__)


class PowerPointParser(BaseParser):
    """
    PowerPoint presentation parser for OpenViking.

    Supports: .pptx

    Converts PowerPoint presentations to Markdown using python-pptx,
    then delegates to MarkdownParser for tree structure creation.
    """

    def __init__(self, config: Optional[ParserConfig] = None, extract_notes: bool = False):
        """
        Initialize PowerPoint parser.

        Args:
            config: Parser configuration
            extract_notes: Whether to extract speaker notes
        """
        from openviking.parse.parsers.markdown import MarkdownParser

        self._md_parser = MarkdownParser(config=config)
        self.config = config or ParserConfig()
        self.extract_notes = extract_notes

    @property
    def supported_extensions(self) -> List[str]:
        return [".pptx"]

    async def parse(self, source: Union[str, Path], instruction: str = "", **kwargs) -> ParseResult:
        """Parse PowerPoint presentation from file path."""
        path = Path(source)

        if path.exists():
            import pptx

            markdown_content = self._convert_to_markdown(path, pptx)
            result = await self._md_parser.parse_content(
                markdown_content, source_path=str(path), instruction=instruction, **kwargs
            )
        else:
            result = await self._md_parser.parse_content(
                str(source), instruction=instruction, **kwargs
            )
        result.source_format = "pptx"
        result.parser_name = "PowerPointParser"
        return result

    async def parse_content(
        self, content: str, source_path: Optional[str] = None, instruction: str = "", **kwargs
    ) -> ParseResult:
        """Parse content - delegates to MarkdownParser."""
        result = await self._md_parser.parse_content(content, source_path, **kwargs)
        result.source_format = "pptx"
        result.parser_name = "PowerPointParser"
        return result

    def _convert_to_markdown(self, path: Path, pptx) -> str:
        """Convert PowerPoint presentation to Markdown string."""
        prs = pptx.Presentation(path)
        markdown_parts = []
        slide_count = len(prs.slides)

        for idx, slide in enumerate(prs.slides, 1):
            slide_parts = []
            slide_parts.append(f"## Slide {idx}/{slide_count}")

            title = self._extract_slide_title(slide)
            if title:
                slide_parts.append(f"### {title}")

            content = self._extract_slide_content(slide)
            if content:
                slide_parts.append(content)

            if self.extract_notes and slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_parts.append(f"**Notes:** {notes}")

            markdown_parts.append("\n\n".join(slide_parts))

        return "\n\n---\n\n".join(markdown_parts)

    def _extract_slide_title(self, slide) -> str:
        """Extract title from a slide."""
        from pptx.enum.shapes import PP_PLACEHOLDER

        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    return shape.text.strip()
        return ""

    def _extract_slide_content(self, slide) -> str:
        """Extract content from slide shapes."""
        from pptx.enum.shapes import PP_PLACEHOLDER

        content_parts = []

        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    continue

            if hasattr(shape, "text") and shape.text.strip():
                if shape.has_table:
                    content_parts.append(self._convert_table(shape.table))
                else:
                    text = shape.text.strip()
                    if text:
                        content_parts.append(text)

        return "\n\n".join(content_parts)

    def _convert_table(self, table) -> str:
        """Convert PowerPoint table to markdown format."""
        if not table.rows:
            return ""

        rows = []
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            rows.append(row_data)

        from openviking.parse.base import format_table_to_markdown

        return format_table_to_markdown(rows, has_header=True)
